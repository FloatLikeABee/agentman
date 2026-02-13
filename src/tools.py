from langchain.tools import Tool
from langchain_community.tools import WikipediaQueryRun
from langchain_community.utilities import WikipediaAPIWrapper
import smtplib
import yfinance as yf
import requests
import json
import logging
from datetime import datetime
from typing import Dict, List, Any, Optional
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse, quote_plus
import time
import random
from .config import settings
from .models import ToolType, ToolConfig, RAGDataInput, DataFormat, LLMProviderType
from .llm_factory import LLMFactory, LLMProvider
from .browser_automation import BrowserAutomationTool


class ToolManager:
    def __init__(self, rag_system=None):
        self.logger = logging.getLogger(__name__)
        self.tools: Dict[str, Tool] = {}
        self.tool_configs: Dict[str, ToolConfig] = {}
        self.rag_system = rag_system
        self._initialize_default_tools()

    def _initialize_default_tools(self):
        """Initialize default tools"""
        # Web Search Tool - Custom implementation with multiple free engines
        search_tool = Tool(
            name="Web Search",
            func=self._web_search,
            description="Search the internet for current information using multiple free search engines. No API key required, unlimited usage."
        )
        self.register_tool(
            "web_search",
            search_tool,
            ToolConfig(
                name="Web Search",
                tool_type=ToolType.WEB_SEARCH,
                description="Search the internet for current information using multiple free search engines. No API key required, unlimited usage.",
                config={}
            )
        )

        # Wikipedia Tool
        wiki_tool = WikipediaQueryRun(api_wrapper=WikipediaAPIWrapper())
        self.register_tool(
            "wikipedia",
            wiki_tool,
            ToolConfig(
                name="Wikipedia Search",
                tool_type=ToolType.WIKIPEDIA,
                description="Search Wikipedia for factual information",
                config={}
            )
        )

        # Calculator Tool
        calculator_tool = Tool(
            name="Calculator",
            func=self._calculator,
            description="Perform mathematical calculations"
        )
        self.register_tool(
            "calculator",
            calculator_tool,
            ToolConfig(
                name="Calculator",
                tool_type=ToolType.CALCULATOR,
                description="Perform mathematical calculations",
                config={}
            )
        )

        # Email Tool
        email_tool = Tool(
            name="Email Sender",
            func=self._send_email,
            description="Send emails using configured SMTP settings"
        )
        self.register_tool(
            "email",
            email_tool,
            ToolConfig(
                name="Email Sender",
                tool_type=ToolType.EMAIL,
                description="Send emails using configured SMTP settings",
                config={
                    "smtp_server": settings.smtp_server,
                    "smtp_port": settings.smtp_port,
                    "smtp_username": settings.smtp_username,
                    "smtp_password": settings.smtp_password
                }
            )
        )

        # Financial Data Tool (using free yfinance API)
        financial_tool = Tool(
            name="Financial Data",
            func=self._get_financial_data,
            description="""Get real-time stock prices and financial data. This is the MOST RELIABLE and RECOMMENDED tool for any stock market queries.

WHEN TO USE:
- User asks about stock prices, stock quotes, or share prices
- User wants financial data, market data, or company valuation
- User asks "what is the price of [STOCK]", "how much is [STOCK] worth", "stock price of [SYMBOL]"
- User wants market cap, volume, 52-week highs/lows, or price changes
- User asks about any publicly traded company's stock

HOW TO USE:
Input format: Extract the stock ticker symbol (1-5 uppercase letters) from the user's query.
Examples:
- "stock price of AAPL" → Use: "AAPL" or "stock price of AAPL"
- "what is MSFT trading at?" → Use: "MSFT" or "MSFT price"
- "TSLA stock quote" → Use: "TSLA" or "TSLA stock"
- "price of Apple stock" → Use: "AAPL" (Apple's ticker is AAPL)

The tool automatically extracts the symbol from phrases like "stock price of", "price of", "quote for", etc.

WHAT IT RETURNS:
- Current stock price (or previous close if market closed)
- Market capitalization
- Trading volume
- Price change and percentage change
- 52-week high and low
- Previous close price

WHY IT'S RELIABLE:
- Uses yfinance library (Yahoo Finance) - industry standard, free, no API key required
- Multiple fallback methods ensure data retrieval even if one method fails
- Supports all major stock exchanges (NYSE, NASDAQ, etc.)
- Handles rate limiting gracefully
- Works for any valid stock ticker symbol

LIMITATIONS:
- Free tier has rate limits (wait 10-30 seconds if rate limited)
- Market data only available during trading hours (shows previous close when closed)
- Requires valid stock ticker symbol (e.g., AAPL, not "Apple Inc.")

IMPORTANT: Always use this tool for stock price queries. It's the best and most reliable option available."""
        )
        self.register_tool(
            "financial",
            financial_tool,
            ToolConfig(
                name="Financial Data",
                tool_type=ToolType.FINANCIAL,
                description="Get real-time stock prices and financial data using yfinance (Yahoo Finance). Most reliable tool for stock market queries. Free, no API key required. Supports all major exchanges. Returns price, market cap, volume, and key metrics.",
                config={}
            )
        )

        # Equalizer Tool
        equalizer_tool = Tool(
            name="Decision Equalizer",
            func=self._equalizer,
            description="Use AI to help make decisions by weighing options given a scenario. Input format: 'scenario:Your decision scenario description'"
        )
        self.register_tool(
            "equalizer",
            equalizer_tool,
            ToolConfig(
                name="Decision Equalizer",
                tool_type=ToolType.EQUALIZER,
                description="Use AI to analyze scenarios and help make weighted decisions",
                config={}
            )
        )

        # Document Reader Tool (PDF, Word, PPT)
        document_reader_tool = Tool(
            name="Document Reader",
            func=self._read_document,
            description="""Analyze documents (PDF, Word, PPT) to summarize, extract key points, answer questions, or find quotes.

INPUT FORMAT: 'action:ACTION_TYPE,query:YOUR_QUERY,url:DOCUMENT_URL' or 'action:ACTION_TYPE,query:YOUR_QUERY,content:DOCUMENT_TEXT'

ACTIONS:
- summarize: Get a comprehensive summary of the document
- extract_key_points: Extract main points and takeaways
- answer_question: Answer specific questions about the document
- find_quotes: Find relevant quotes on a topic

EXAMPLES:
- 'action:summarize,url:https://example.com/doc.pdf'
- 'action:extract_key_points,content:Your document text here...'
- 'action:answer_question,query:What is the main conclusion?,url:https://example.com/report.pdf'
- 'action:find_quotes,query:innovation,content:Document text about innovation...'

Supports: PDF, DOCX, PPTX, TXT files via URL or direct text content."""
        )
        self.register_tool(
            "document_reader",
            document_reader_tool,
            ToolConfig(
                name="Document Reader",
                tool_type=ToolType.DOCUMENT_READER,
                description="Analyze documents (PDF, Word, PPT) - summarize, extract key points, answer questions, find quotes",
                config={}
            )
        )

        # YouTube/Video Summarizer Tool
        youtube_summarizer_tool = Tool(
            name="YouTube Summarizer",
            func=self._summarize_youtube,
            description="""Analyze YouTube videos to get summaries, timestamps, key insights, and action items.

INPUT FORMAT: 'url:YOUTUBE_URL' or 'url:YOUTUBE_URL,action:ACTION_TYPE'

ACTIONS (optional, default is 'full'):
- full: Complete analysis with summary, timestamps, insights, and action items
- summary: Just the summary
- timestamps: Key moments with timestamps
- insights: Main insights and learnings
- action_items: Actionable takeaways

EXAMPLES:
- 'url:https://www.youtube.com/watch?v=VIDEO_ID'
- 'url:https://youtu.be/VIDEO_ID,action:summary'
- 'url:https://www.youtube.com/watch?v=VIDEO_ID,action:timestamps'

Works with any public YouTube video."""
        )
        self.register_tool(
            "youtube_summarizer",
            youtube_summarizer_tool,
            ToolConfig(
                name="YouTube Summarizer",
                tool_type=ToolType.YOUTUBE_SUMMARIZER,
                description="Analyze YouTube videos - get summaries, timestamps, key insights, and action items",
                config={}
            )
        )

        # Academic Paper Search Tool
        academic_search_tool = Tool(
            name="Academic Search",
            func=self._search_academic,
            description="""Search and analyze academic papers from multiple sources (Semantic Scholar, arXiv, CrossRef).

INPUT FORMAT: 'query:SEARCH_QUERY' or 'query:SEARCH_QUERY,action:ACTION_TYPE'

ACTIONS (optional, default is 'search'):
- search: Find papers matching the query
- summarize: Search and summarize top papers
- explain_simple: Explain papers in simple terms (ELI5)
- extract_methods: Extract methodologies and formulas

EXAMPLES:
- 'query:machine learning transformers'
- 'query:quantum computing applications,action:summarize'
- 'query:CRISPR gene editing,action:explain_simple'
- 'query:neural network optimization,action:extract_methods'

Returns: Paper titles, authors, abstracts, citations, and download links when available."""
        )
        self.register_tool(
            "academic_search",
            academic_search_tool,
            ToolConfig(
                name="Academic Search",
                tool_type=ToolType.ACADEMIC_SEARCH,
                description="Search academic papers - find, summarize, explain simply, extract methods & formulas",
                config={}
            )
        )

        # Mind Map Generator Tool
        mind_map_tool = Tool(
            name="Mind Map Generator",
            func=self._generate_mind_map,
            description="""Generate a structured mind map from any topic for learning and brainstorming.

INPUT FORMAT: 'topic:YOUR_TOPIC' or 'topic:YOUR_TOPIC,depth:DEPTH_LEVEL'

OPTIONS:
- topic: The main subject to map (required)
- depth: How deep to go (1-3, default 2)
  - 1: High-level overview
  - 2: Moderate detail (recommended)
  - 3: Deep dive with sub-concepts

EXAMPLES:
- 'topic:Machine Learning'
- 'topic:Climate Change,depth:3'
- 'topic:Starting a Business'
- 'topic:Quantum Physics,depth:1'

OUTPUT INCLUDES:
- Central concept and definition
- Main branches/categories
- Sub-branches with relationships
- Connections between concepts
- Key terms and definitions

Perfect for learning new subjects, brainstorming, and organizing ideas."""
        )
        self.register_tool(
            "mind_map",
            mind_map_tool,
            ToolConfig(
                name="Mind Map Generator",
                tool_type=ToolType.MIND_MAP,
                description="Generate structured mind maps for any topic - concepts, relationships, hierarchies",
                config={}
            )
        )

        # Debate / Pros-Cons Analyzer Tool
        debate_analyzer_tool = Tool(
            name="Debate Analyzer",
            func=self._analyze_debate,
            description="""Analyze decisions or topics with comprehensive pros/cons analysis.

INPUT FORMAT: 'question:YOUR_QUESTION' or 'topic:YOUR_TOPIC'

EXAMPLES:
- 'question:Should I quit my job?'
- 'question:Should I start a business or get an MBA?'
- 'topic:Remote work vs office work'
- 'question:Is it worth buying a house in 2024?'

OUTPUT INCLUDES:
- Pros (advantages, benefits)
- Cons (disadvantages, drawbacks)
- Risks (potential negative outcomes)
- Opportunity costs (what you give up)
- Long-term impact analysis
- Key considerations
- Balanced recommendation

Great for major life decisions, business choices, and strategic planning."""
        )
        self.register_tool(
            "debate_analyzer",
            debate_analyzer_tool,
            ToolConfig(
                name="Debate Analyzer",
                tool_type=ToolType.DEBATE_ANALYZER,
                description="Analyze decisions with pros, cons, risks, opportunity costs, and long-term impact",
                config={}
            )
        )

        # First-Principles Reasoner Tool
        first_principles_tool = Tool(
            name="First Principles Reasoner",
            func=self._first_principles_reasoning,
            description="""Break down complex problems to fundamental truths using first-principles thinking.

INPUT FORMAT: 'problem:YOUR_PROBLEM' or 'idea:YOUR_IDEA'

EXAMPLES:
- 'problem:How can we make electric cars more affordable?'
- 'idea:A subscription service for home-cooked meals'
- 'problem:Why do startups fail?'
- 'idea:Decentralized social media platform'

OUTPUT INCLUDES:
1. DECOMPOSITION: Break down to fundamental components
2. ASSUMPTIONS: Identify and question assumptions
3. FUNDAMENTAL TRUTHS: Core facts that can't be reduced further
4. REBUILD: Reconstruct from first principles
5. NOVEL INSIGHTS: New perspectives from this analysis
6. ACTION STEPS: Practical next steps

Used by Elon Musk, scientists, and innovators. Great for startups, physics problems, business logic, and creative problem-solving."""
        )
        self.register_tool(
            "first_principles",
            first_principles_tool,
            ToolConfig(
                name="First Principles Reasoner",
                tool_type=ToolType.FIRST_PRINCIPLES,
                description="Break ideas down to fundamentals - great for startups, physics, business logic",
                config={}
            )
        )

        # Image Generator Tool (using Pollinations API)
        image_generator_tool = Tool(
            name="Image Generator",
            func=self._generate_image,
            description="""Generate images via Pollinations API. Returns an image URL - DO NOT describe or recreate the image, just return the URL to the user.

INPUT: Just the image description text, e.g. 'a cute cat' or 'futuristic city at sunset'

IMPORTANT: This tool calls an external API and returns an image URL. Simply pass the URL result back to the user. Do not attempt to describe, draw, or recreate the image yourself."""
        )
        self.register_tool(
            "image_generator",
            image_generator_tool,
            ToolConfig(
                name="Image Generator",
                tool_type=ToolType.IMAGE_GENERATOR,
                description="Generate images from text descriptions using AI (Pollinations Flux model)",
                config={}
            )
        )

        # Story / Script Generator Tool
        story_generator_tool = Tool(
            name="Story Generator",
            func=self._generate_story,
            description="""Generate creative stories, scripts, and narratives for various media formats.

INPUT FORMAT: 'type:FORMAT,premise:YOUR_PREMISE' or structured input

SUPPORTED FORMATS:
- comic: Comic book scripts with panels, dialogue, and visual descriptions
- game: Game narratives with dialogue trees, quests, and character arcs
- movie: Screenplay format with scenes, action, and dialogue
- ad: Advertising scripts for commercials, social media, or print
- short_story: Traditional narrative fiction
- podcast: Audio script with host/guest dialogue

EXAMPLES:
- 'type:comic,premise:A superhero who can only use powers when singing'
- 'type:game,premise:Post-apocalyptic survival RPG,genre:sci-fi,tone:dark'
- 'type:movie,premise:Two rival chefs fall in love,genre:romantic comedy,length:short'
- 'type:ad,premise:Eco-friendly water bottle,target:millennials,duration:30sec'
- 'type:short_story,premise:A time traveler stuck in ancient Rome,style:humorous'

OPTIONS:
- type: Format type (required)
- premise: Story concept or idea (required)
- genre: Genre (action, comedy, drama, horror, sci-fi, fantasy, romance, thriller)
- tone: Emotional tone (dark, light, humorous, serious, inspirational)
- length: short, medium, long (affects detail level)
- characters: Key character descriptions
- setting: Time and place
- target: Target audience (for ads)

OUTPUT INCLUDES:
- Title and logline
- Character profiles
- Scene/panel breakdowns
- Dialogue
- Visual/action descriptions
- Format-specific elements (panels for comics, quest structure for games, etc.)"""
        )
        self.register_tool(
            "story_generator",
            story_generator_tool,
            ToolConfig(
                name="Story Generator",
                tool_type=ToolType.STORY_GENERATOR,
                description="Generate stories and scripts for comics, games, movies, and ads",
                config={}
            )
        )

        # Task Planner Tool
        task_planner_tool = Tool(
            name="Task Planner",
            func=self._plan_tasks,
            description="""Transform goals into actionable task plans with steps, timelines, and checklists.

INPUT FORMAT: 'goal:YOUR_GOAL' or structured input with options

EXAMPLES:
- 'goal:Launch a mobile app for fitness tracking'
- 'goal:Plan a wedding for 100 guests,budget:$20000,deadline:6 months'
- 'goal:Learn Python programming from scratch,hours_per_week:10'
- 'goal:Renovate kitchen,constraints:must keep working during renovation'
- 'goal:Start a YouTube channel about cooking'

OPTIONS:
- goal: What you want to achieve (required)
- deadline: Target completion date/timeframe
- budget: Available budget
- hours_per_week: Time commitment available
- constraints: Any limitations or requirements
- priority: What matters most (speed, quality, cost)
- team_size: Number of people available
- skill_level: beginner, intermediate, advanced

OUTPUT INCLUDES:
1. GOAL ANALYSIS: Clarified objective and success criteria
2. PHASES: Major project phases with milestones
3. DETAILED STEPS: Actionable tasks for each phase
4. TIMELINE: Suggested schedule with dependencies
5. CHECKLIST: Trackable to-do items with checkboxes
6. RESOURCES: Tools, skills, or help needed
7. RISKS: Potential obstacles and mitigation strategies
8. QUICK WINS: Easy early tasks to build momentum

Perfect for project planning, learning goals, life events, business launches, and personal development."""
        )
        self.register_tool(
            "task_planner",
            task_planner_tool,
            ToolConfig(
                name="Task Planner",
                tool_type=ToolType.TASK_PLANNER,
                description="Transform goals into actionable plans with steps, timelines, and checklists",
                config={}
            )
        )

        # Multi-Agent Tool
        multi_agent_tool = Tool(
            name="Multi-Agent Collaborator",
            func=self._multi_agent_process,
            description="""Simulate multiple AI agents collaborating on a task, each with a specialized role.

INPUT FORMAT: 'task:YOUR_TASK' or structured input with options

AGENT ROLES:
- Researcher: Gathers information, finds facts, explores the topic deeply
- Writer: Creates content, drafts text, structures information
- Critic: Evaluates quality, finds flaws, suggests improvements
- Summarizer: Distills key points, creates concise versions

EXAMPLES:
- 'task:Write a blog post about sustainable fashion'
- 'task:Analyze the pros and cons of remote work,depth:thorough'
- 'task:Create a business proposal for a coffee subscription service'
- 'task:Explain quantum computing for beginners,style:engaging'
- 'task:Review and improve this essay about climate change'

OPTIONS:
- task: The task to collaborate on (required)
- depth: quick, standard, thorough (affects how many rounds)
- focus: research, writing, critique, or balanced
- style: formal, casual, academic, creative
- iterations: Number of improvement rounds (1-3)

PROCESS:
1. RESEARCHER explores the topic, gathers key information
2. WRITER creates initial content based on research
3. CRITIC reviews and identifies improvements
4. WRITER revises based on feedback
5. SUMMARIZER creates final polished version

OUTPUT INCLUDES:
- Research findings from the Researcher
- Initial draft from the Writer
- Critique and suggestions from the Critic
- Revised content addressing feedback
- Final summary with key takeaways

This simulates a collaborative team process for higher quality outputs."""
        )
        self.register_tool(
            "multi_agent",
            multi_agent_tool,
            ToolConfig(
                name="Multi-Agent Collaborator",
                tool_type=ToolType.MULTI_AGENT,
                description="Multiple AI agents collaborate: researcher, writer, critic, and summarizer",
                config={}
            )
        )

        # Browser Automation Tool
        browser_automation = BrowserAutomationTool()
        browser_tool = Tool(
            name="Browser Automation",
            func=browser_automation.execute,
            description="""Control a web browser to perform tasks automatically. Uses AI agent with Playwright to follow instructions.

INPUT FORMAT: Natural language instructions describing what to do in the browser.

CAPABILITIES:
- Navigate to websites
- Click buttons and links
- Fill forms
- Extract text and data
- Take screenshots
- Scroll pages
- Wait for elements
- Select dropdown options

EXAMPLES:
- "Go to google.com and search for 'Python programming'"
- "Navigate to example.com, fill the contact form with name 'John Doe' and email 'john@example.com', then submit"
- "Open github.com, search for 'langchain', and get the first 5 repository names"
- "Go to news.ycombinator.com, scroll down, and take a screenshot"
- "Navigate to amazon.com, search for 'laptop', and get the price of the first result"

The agent will automatically break down your instructions into browser actions and execute them step by step.
Returns: Summary of actions taken, final URL, and page content summary."""
        )
        self.register_tool(
            "browser_automation",
            browser_tool,
            ToolConfig(
                name="Browser Automation",
                tool_type=ToolType.BROWSER_AUTOMATION,
                description="Control a web browser using AI agent with Playwright to follow natural language instructions",
                config={}
            )
        )

    def register_tool(self, tool_id: str, tool: Tool, config: ToolConfig):
        """Register a new tool"""
        self.tools[tool_id] = tool
        self.tool_configs[tool_id] = config
        self.logger.info(f"Registered tool: {tool_id}")

    def get_tool(self, tool_id: str) -> Optional[Tool]:
        """Get a tool by ID or name (case-insensitive)"""
        # First try exact match
        if tool_id in self.tools:
            return self.tools.get(tool_id)
        
        # Try case-insensitive match
        tool_id_lower = tool_id.lower()
        for key in self.tools.keys():
            if key.lower() == tool_id_lower:
                return self.tools.get(key)
        
        # Try matching by tool config name
        for key, config in self.tool_configs.items():
            if config.name.lower() == tool_id_lower or config.name == tool_id:
                return self.tools.get(key)
        
        return None

    def list_tools(self) -> List[Dict[str, Any]]:
        """List all available tools"""
        return [
            {
                'id': tool_id,
                'name': config.name,
                'tool_type': config.tool_type,
                'description': config.description,
                'is_active': config.is_active,
                'config': config.config
            }
            for tool_id, config in self.tool_configs.items()
        ]

    def update_tool_config(self, tool_id: str, config: ToolConfig) -> bool:
        """Update tool configuration"""
        try:
            if tool_id in self.tool_configs:
                self.tool_configs[tool_id] = config
                self.logger.info(f"Updated tool config: {tool_id}")
                return True
            return False
        except Exception as e:
            self.logger.error(f"Error updating tool config {tool_id}: {e}")
            return False

    def _calculator(self, expression: str) -> str:
        """Calculator tool function"""
        try:
            # Safe evaluation of mathematical expressions
            allowed_chars = set('0123456789+-*/(). ')
            if not all(c in allowed_chars for c in expression):
                return "Error: Invalid characters in expression"
            
            result = eval(expression)
            return str(result)
        except Exception as e:
            return f"Error in calculation: {str(e)}"

    def _send_email(self, email_data: str) -> str:
        """Email sending tool function"""
        try:
            # Parse email data (expected format: "to:email@example.com,subject:Subject,body:Email body")
            data_parts = email_data.split(',')
            email_dict = {}
            for part in data_parts:
                if ':' in part:
                    key, value = part.split(':', 1)
                    email_dict[key.strip()] = value.strip()

            to_email = email_dict.get('to', '')
            subject = email_dict.get('subject', '')
            body = email_dict.get('body', '')

            if not all([to_email, subject, body]):
                return "Error: Missing required email fields (to, subject, body)"

            # Check if SMTP is configured
            if not settings.smtp_server or not settings.smtp_username:
                return "Error: SMTP not configured"

            # Create message
            msg = MIMEMultipart()
            msg['From'] = settings.smtp_username
            msg['To'] = to_email
            msg['Subject'] = subject
            msg.attach(MIMEText(body, 'plain'))

            # Send email
            server = smtplib.SMTP(settings.smtp_server, settings.smtp_port)
            server.starttls()
            server.login(settings.smtp_username, settings.smtp_password)
            text = msg.as_string()
            server.sendmail(settings.smtp_username, to_email, text)
            server.quit()

            return f"Email sent successfully to {to_email}"

        except Exception as e:
            return f"Error sending email: {str(e)}"

    def _web_search(self, query: str) -> str:
        """Custom web search using multiple free search engines with no limits"""
        try:
            results = []
            query_encoded = quote_plus(query)
            
            # Try Tavily API first if API key is available (optional, has free tier)
            tavily_api_key = getattr(settings, 'tavily_api_key', None)
            if tavily_api_key:
                try:
                    self.logger.info("Trying Tavily API for web search")
                    tavily_url = "https://api.tavily.com/search"
                    payload = {
                        "api_key": tavily_api_key,
                        "query": query,
                        "search_depth": "basic",
                        "max_results": 10
                    }
                    response = requests.post(tavily_url, json=payload, timeout=10)
                    if response.status_code == 200:
                        data = response.json()
                        if data.get('results'):
                            for result in data['results']:
                                results.append({
                                    'title': result.get('title', 'No title'),
                                    'url': result.get('url', ''),
                                    'snippet': result.get('content', 'No description available')
                                })
                            if results:
                                self.logger.info(f"Successfully retrieved {len(results)} results from Tavily")
                                formatted_results = []
                                for i, result in enumerate(results[:10], 1):
                                    formatted_results.append(
                                        f"{i}. {result['title']}\n   URL: {result['url']}\n   {result['snippet']}"
                                    )
                                return "\n\n".join(formatted_results)
                except Exception as e:
                    self.logger.debug(f"Tavily API failed: {e}, falling back to free search engines")
            
            # Fallback to free search engines (no API key required, unlimited)
            # Try multiple free search engines as fallbacks
            search_engines = [
                {
                    "name": "Startpage",
                    "url": f"https://www.startpage.com/sp/search?query={query_encoded}",
                    "selectors": {
                        "results": "div.w-gl__result",
                        "title": "h3 a",
                        "link": "h3 a",
                        "snippet": "p.w-gl__description"
                    }
                },
                {
                    "name": "Qwant",
                    "url": f"https://www.qwant.com/?q={query_encoded}&t=web",
                    "selectors": {
                        "results": "div.web-result",
                        "title": "a",
                        "link": "a",
                        "snippet": "p.web-result-description"
                    }
                },
                {
                    "name": "Ecosia",
                    "url": f"https://www.ecosia.org/search?q={query_encoded}",
                    "selectors": {
                        "results": "div.result",
                        "title": "h2 a",
                        "link": "h2 a",
                        "snippet": "p.result-snippet"
                    }
                }
            ]
            
            # Headers to mimic a real browser
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.9',
                'Accept-Encoding': 'gzip, deflate, br',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
            }
            
            # Try each search engine until we get results
            for engine in search_engines:
                try:
                    self.logger.info(f"Trying {engine['name']} for query: {query}")
                    response = requests.get(engine['url'], headers=headers, timeout=10)
                    response.raise_for_status()
                    
                    soup = BeautifulSoup(response.content, 'html.parser')
                    result_elements = soup.select(engine['selectors']['results'])
                    
                    if result_elements:
                        for element in result_elements[:5]:  # Limit to top 5 results per engine
                            try:
                                title_elem = element.select_one(engine['selectors']['title'])
                                link_elem = element.select_one(engine['selectors']['link'])
                                snippet_elem = element.select_one(engine['selectors']['snippet'])
                                
                                if title_elem and link_elem:
                                    title = title_elem.get_text(strip=True)
                                    link = link_elem.get('href', '')
                                    
                                    # Handle relative URLs
                                    if link and not link.startswith('http'):
                                        link = urljoin(engine['url'], link)
                                    
                                    snippet = snippet_elem.get_text(strip=True) if snippet_elem else "No description available"
                                    
                                    if title and link:  # Only add if we have both
                                        results.append({
                                            'title': title,
                                            'url': link,
                                            'snippet': snippet
                                        })
                            except Exception as e:
                                self.logger.debug(f"Error parsing result element: {e}")
                                continue
                        
                        if results:
                            self.logger.info(f"Successfully retrieved {len(results)} results from {engine['name']}")
                            break  # We got results, no need to try other engines
                            
                except requests.exceptions.RequestException as e:
                    self.logger.warning(f"Failed to search with {engine['name']}: {e}")
                    continue
                except Exception as e:
                    self.logger.warning(f"Error parsing {engine['name']} results: {e}")
                    continue
            
            # If no results from scraping, try DuckDuckGo instant answer API as final fallback
            if not results:
                try:
                    self.logger.info("Trying DuckDuckGo instant answer API as final fallback")
                    ddg_url = f"https://api.duckduckgo.com/?q={query_encoded}&format=json&no_html=1&skip_disambig=1"
                    response = requests.get(ddg_url, headers=headers, timeout=10)
                    if response.status_code == 200:
                        data = response.json()
                        if data.get('AbstractText'):
                            results.append({
                                'title': data.get('Heading', 'DuckDuckGo Result'),
                                'url': data.get('AbstractURL', ''),
                                'snippet': data.get('AbstractText', '')
                            })
                except Exception as e:
                    self.logger.debug(f"DuckDuckGo API fallback failed: {e}")
            
            # Format results
            if results:
                formatted_results = []
                for i, result in enumerate(results[:10], 1):  # Limit to top 10 total results
                    formatted_results.append(
                        f"{i}. {result['title']}\n   URL: {result['url']}\n   {result['snippet']}"
                    )
                return "\n\n".join(formatted_results)
            else:
                return f"No search results found for: {query}. Please try rephrasing your query."
                
        except Exception as e:
            self.logger.error(f"Error in web search: {e}")
            return f"Error performing web search: {str(e)}. Please try again."

    def _get_financial_data(self, query: str) -> str:
        """Financial data tool function using free yfinance API (no API key required)"""
        try:
            import time
            import re
            from datetime import datetime
            
            # Parse query for stock symbol - improved extraction
            symbol = None
            
            # Try to extract symbol from common patterns
            words = query.split()
            for i, word in enumerate(words):
                # Look for uppercase ticker symbols (1-5 characters)
                if len(word) <= 5 and word.isupper() and word.isalpha():
                    symbol = word
                    break
                # Look for patterns like "AAPL stock" or "stock AAPL"
                if word.lower() in ['stock', 'price', 'quote', 'of'] and i + 1 < len(words):
                    potential_symbol = words[i + 1].upper().strip('.,;:!?')
                    if len(potential_symbol) <= 5 and potential_symbol.isalpha():
                        symbol = potential_symbol
                        break
            
            if not symbol:
                # Try to find symbol in the query more flexibly using regex
                symbol_match = re.search(r'\b([A-Z]{1,5})\b', query.upper())
                if symbol_match:
                    symbol = symbol_match.group(1)
            
            if not symbol:
                return "FORMAT: Please provide a stock symbol (e.g., AAPL, MSFT, TSLA). Example: 'stock price of AAPL' or 'AAPL price'"
            
            # Get stock data using yfinance (free, no API key required)
            # Use history first as it's less likely to hit rate limits
            try:
                stock = yf.Ticker(symbol)
                current_price = None
                info = {}
                
                # Method 1: Try history first (less rate-limited than info)
                try:
                    # Try 1 day history with daily intervals (most reliable)
                    hist = stock.history(period="1d", interval="1d")
                    if not hist.empty and len(hist) > 0:
                        current_price = float(hist['Close'].iloc[-1])
                except:
                    try:
                        # Fallback: Try 5 day history
                        hist = stock.history(period="5d", interval="1d")
                        if not hist.empty and len(hist) > 0:
                            current_price = float(hist['Close'].iloc[-1])
                    except:
                        pass
                
                # Method 2: Try info if history worked (for additional data)
                # Only try info if we got price from history, to avoid rate limits
                if current_price is not None:
                    try:
                        info = stock.info
                        if info and not info.get('currentPrice'):
                            # If info doesn't have currentPrice, use the one from history
                            pass
                    except Exception as e1:
                        # If info fails due to rate limit, that's okay - we have price from history
                        self.logger.debug(f"Info unavailable for {symbol} (may be rate limited): {e1}")
                        info = {}
                
                # Method 3: If history failed, try info as last resort
                if current_price is None:
                    try:
                        info = stock.info
                        if info:
                            # Try various price fields
                            current_price = (info.get('currentPrice') or 
                                            info.get('regularMarketPrice') or 
                                            info.get('regularMarketPreviousClose') or
                                            info.get('previousClose') or
                                            info.get('ask') or
                                            info.get('bid'))
                            if current_price:
                                current_price = float(current_price)
                    except Exception as e1:
                        # Check if it's a rate limit error
                        error_str = str(e1).lower()
                        if 'rate limit' in error_str or 'too many' in error_str:
                            # Return helpful message about rate limiting
                            return f"Stock: {symbol} | Status: Rate limited by Yahoo Finance. Please wait 30-60 seconds and try again. The free API has usage limits."
                        self.logger.debug(f"Error getting info for {symbol}: {e1}")
                        info = {}
                
                # Method 4: Try fast_info as final fallback
                if current_price is None:
                    try:
                        fast_info = stock.fast_info
                        if hasattr(fast_info, 'last_price') and fast_info.last_price:
                            current_price = float(fast_info.last_price)
                        elif hasattr(fast_info, 'regular_market_price') and fast_info.regular_market_price:
                            current_price = float(fast_info.regular_market_price)
                    except:
                        pass
                
                # Build structured response that won't break agent parsing
                result_parts = [f"Stock: {symbol}"]
                
                if current_price:
                    result_parts.append(f"Price: ${current_price:.2f}")
                elif info and info.get('previousClose'):
                    # If we have previous close but no current price, show that
                    prev_close = float(info.get('previousClose'))
                    result_parts.append(f"Previous Close: ${prev_close:.2f}")
                    result_parts.append("Note: Current price unavailable (market may be closed)")
                else:
                    # Last resort: try to get any price data
                    if info:
                        for price_key in ['regularMarketPreviousClose', 'previousClose', 'fiftyTwoWeekHigh', 'fiftyTwoWeekLow']:
                            if price_key in info and info[price_key]:
                                price_val = float(info[price_key])
                                result_parts.append(f"{price_key.replace('fiftyTwo', '52W').replace('regularMarket', '')}: ${price_val:.2f}")
                                break
                    
                    if len(result_parts) == 1:  # Only has "Stock: SYMBOL"
                        result_parts.append("Price: Unable to fetch current price. Market may be closed or symbol invalid.")
                
                # Add more data if available
                if info:
                    if 'marketCap' in info and info['marketCap']:
                        try:
                            market_cap = float(info['marketCap'])
                            if market_cap >= 1e12:
                                market_cap_str = f"${market_cap/1e12:.2f}T"
                            elif market_cap >= 1e9:
                                market_cap_str = f"${market_cap/1e9:.2f}B"
                            elif market_cap >= 1e6:
                                market_cap_str = f"${market_cap/1e6:.2f}M"
                            else:
                                market_cap_str = f"${market_cap:,.0f}"
                            result_parts.append(f"Market Cap: {market_cap_str}")
                        except:
                            pass
                    
                    if 'volume' in info and info['volume']:
                        try:
                            result_parts.append(f"Volume: {int(info['volume']):,}")
                        except:
                            pass
                    
                    if current_price and 'previousClose' in info and info['previousClose']:
                        try:
                            prev_close = float(info['previousClose'])
                            change = current_price - prev_close
                            change_pct = (change / prev_close) * 100
                            result_parts.append(f"Previous Close: ${prev_close:.2f}")
                            result_parts.append(f"Change: ${change:+.2f} ({change_pct:+.2f}%)")
                        except:
                            pass
                    
                    if '52WeekHigh' in info and info['52WeekHigh']:
                        try:
                            result_parts.append(f"52W High: ${float(info['52WeekHigh']):.2f}")
                        except:
                            pass
                    
                    if '52WeekLow' in info and info['52WeekLow']:
                        try:
                            result_parts.append(f"52W Low: ${float(info['52WeekLow']):.2f}")
                        except:
                            pass
                
                return "\n".join(result_parts)
                
            except Exception as yf_error:
                # Handle yfinance-specific errors gracefully with structured format
                error_msg = str(yf_error).lower()
                if 'rate limit' in error_msg or 'too many' in error_msg or '429' in str(yf_error):
                    return f"Stock: {symbol} | Status: Rate limited. Please wait 10 seconds and try again."
                elif 'not found' in error_msg or 'invalid' in error_msg or '404' in str(yf_error):
                    return f"Stock: {symbol} | Status: Symbol not found. Please verify the stock symbol is correct."
                else:
                    # Try one more time with a simpler approach
                    try:
                        stock = yf.Ticker(symbol)
                        fast_info = stock.fast_info
                        if hasattr(fast_info, 'last_price'):
                            price = float(fast_info.last_price)
                            return f"Stock: {symbol}\nPrice: ${price:.2f}"
                    except:
                        return f"Stock: {symbol} | Status: Unable to fetch data. Please verify symbol and try again later."

        except Exception as e:
            # Return structured error format that won't break agent parsing
            return f"Status: Error retrieving financial data. Please provide a valid stock symbol (e.g., AAPL, MSFT)."

    def create_custom_tool(self, tool_id: str, name: str, description: str, func) -> bool:
        """Create a custom tool"""
        try:
            custom_tool = Tool(
                name=name,
                func=func,
                description=description
            )
            
            config = ToolConfig(
                name=name,
                tool_type=ToolType.CUSTOM,
                description=description,
                config={}
            )
            
            self.register_tool(tool_id, custom_tool, config)
            return True
            
        except Exception as e:
            self.logger.error(f"Error creating custom tool {tool_id}: {e}")
            return False

    def _equalizer(self, input_data: str) -> str:
        """Equalizer tool: Use AI to help make decisions by weighing options"""
        try:
            # Parse input: "scenario:Your decision scenario description"
            if ':' in input_data:
                key, scenario = input_data.split(':', 1)
                scenario = scenario.strip()
            else:
                scenario = input_data.strip()
            
            if not scenario:
                return "Error: Please provide a decision scenario. Format: 'scenario:Your scenario description'"
            
            # Create LLM caller
            try:
                # Determine provider from settings (case-insensitive)
                provider_str = settings.default_llm_provider.lower().strip()
                self.logger.info(f"Tool requested provider: '{settings.default_llm_provider}' (normalized: '{provider_str}')")
                
                if provider_str == "gemini":
                    provider = LLMProvider.GEMINI
                    api_key = settings.gemini_api_key
                    model = settings.gemini_default_model
                    self.logger.info("Selected Gemini provider for equalizer tool")
                elif provider_str == "qwen":
                    provider = LLMProvider.QWEN
                    api_key = settings.qwen_api_key
                    model = settings.qwen_default_model
                    self.logger.info("Selected Qwen provider for equalizer tool")
                else:
                    # Fallback: try Qwen first, then Gemini
                    self.logger.warning(f"Unknown provider '{settings.default_llm_provider}', trying Qwen first")
                    try:
                        provider = LLMProvider.QWEN
                        api_key = settings.qwen_api_key
                        model = settings.qwen_default_model
                        llm_caller = LLMFactory.create_caller(provider=provider, api_key=api_key, model=model)
                        self.logger.info("Successfully initialized Qwen as fallback")
                    except Exception as e:
                        self.logger.warning(f"Failed to initialize Qwen: {e}, trying Gemini")
                        provider = LLMProvider.GEMINI
                        api_key = settings.gemini_api_key
                        model = settings.gemini_default_model
                        llm_caller = LLMFactory.create_caller(provider=provider, api_key=api_key, model=model)
                        self.logger.info("Successfully initialized Gemini as fallback")
                
                if 'llm_caller' not in locals():
                    self.logger.info(f"Initializing {provider.value} with model {model}")
                    llm_caller = LLMFactory.create_caller(
                        provider=provider,
                        api_key=api_key,
                        model=model
                    )
                
                # Create decision analysis prompt
                decision_prompt = f"""You are a decision-making assistant. Analyze the following scenario and help make a well-reasoned decision by weighing the options.

Scenario:
{scenario}

Please provide:
1. Identify all possible options/choices
2. List pros and cons for each option
3. Assign weights/scores to each factor (1-10 scale)
4. Calculate weighted scores for each option
5. Provide a clear recommendation with reasoning
6. Include any risks or considerations

Format your response as structured JSON:
{{
    "options": [
        {{
            "name": "Option 1",
            "pros": ["pro1", "pro2"],
            "cons": ["con1", "con2"],
            "factors": [
                {{"factor": "Factor name", "weight": 8, "score": 7}},
                ...
            ],
            "total_score": 0,
            "weighted_score": 0
        }},
        ...
    ],
    "analysis": "Overall analysis of the decision",
    "recommendation": {{
        "option": "Recommended option",
        "reasoning": "Why this option is recommended",
        "confidence": "high/medium/low"
    }},
    "risks": ["risk1", "risk2"],
    "considerations": ["consideration1", "consideration2"]
}}

Calculate weighted scores by: sum(factor.weight * factor.score) / sum(factor.weight) for each option."""
                
                # Call AI for decision analysis
                self.logger.info("Calling AI for decision analysis...")
                response = llm_caller.generate(decision_prompt)
                
                # Try to extract JSON from response
                try:
                    # Find JSON in the response
                    if '```json' in response:
                        response = response.split('```json')[1].split('```')[0].strip()
                    elif '```' in response:
                        response = response.split('```')[1].split('```')[0].strip()
                    
                    decision_data = json.loads(response)
                    
                    # Format the response nicely
                    result = f"Decision Analysis for: {scenario}\n\n"
                    result += f"Analysis: {decision_data.get('analysis', 'N/A')}\n\n"
                    
                    result += "Options Evaluation:\n"
                    for i, option in enumerate(decision_data.get('options', []), 1):
                        result += f"\n{i}. {option.get('name', 'Option')}\n"
                        result += f"   Pros: {', '.join(option.get('pros', []))}\n"
                        result += f"   Cons: {', '.join(option.get('cons', []))}\n"
                        result += f"   Weighted Score: {option.get('weighted_score', 0):.2f}\n"
                    
                    recommendation = decision_data.get('recommendation', {})
                    result += f"\nRecommendation: {recommendation.get('option', 'N/A')}\n"
                    result += f"Reasoning: {recommendation.get('reasoning', 'N/A')}\n"
                    result += f"Confidence: {recommendation.get('confidence', 'N/A')}\n"
                    
                    if decision_data.get('risks'):
                        result += f"\nRisks: {', '.join(decision_data.get('risks', []))}\n"
                    if decision_data.get('considerations'):
                        result += f"Considerations: {', '.join(decision_data.get('considerations', []))}\n"
                    
                    return result
                    
                except json.JSONDecodeError:
                    # If JSON parsing fails, return the raw response
                    return f"Decision Analysis:\n\n{response}"
                    
            except Exception as e:
                self.logger.error(f"Error calling AI for decision analysis: {e}")
                return f"Error in AI decision analysis: {str(e)}"
                
        except Exception as e:
            self.logger.error(f"Error in equalizer tool: {e}")
            return f"Error in equalizer: {str(e)}"

    def _get_llm_caller(self):
        """Helper to get an LLM caller based on settings"""
        provider_str = settings.default_llm_provider.lower().strip()
        
        if provider_str == "gemini":
            provider = LLMProvider.GEMINI
            api_key = settings.gemini_api_key
            model = settings.gemini_default_model
        elif provider_str == "qwen":
            provider = LLMProvider.QWEN
            api_key = settings.qwen_api_key
            model = settings.qwen_default_model
        elif provider_str == "mistral":
            provider = LLMProvider.MISTRAL
            api_key = settings.mistral_api_key
            model = settings.mistral_default_model
        elif provider_str == "groq":
            provider = LLMProvider.GROQ
            api_key = getattr(settings, "groq_api_key", "")
            model = getattr(settings, "groq_default_model", "llama-3.3-70b-versatile")
        else:
            # Default to Qwen
            provider = LLMProvider.QWEN
            api_key = settings.qwen_api_key
            model = settings.qwen_default_model
        
        return LLMFactory.create_caller(provider=provider, api_key=api_key, model=model)

    def _read_document(self, input_data: str) -> str:
        """Document Reader: Analyze PDFs, Word docs, PPTs"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            action = params.get('action', 'summarize')
            query = params.get('query', '')
            url = params.get('url', '')
            content = params.get('content', '')
            
            # Get document content
            document_text = content
            
            if url and not content:
                try:
                    # Fetch document from URL
                    headers = {
                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                    }
                    response = requests.get(url, headers=headers, timeout=30)
                    response.raise_for_status()
                    
                    content_type = response.headers.get('content-type', '').lower()
                    
                    if 'pdf' in content_type or url.lower().endswith('.pdf'):
                        # Extract text from PDF
                        try:
                            import io
                            try:
                                import pypdf
                                pdf_reader = pypdf.PdfReader(io.BytesIO(response.content))
                                document_text = "\n".join([page.extract_text() or '' for page in pdf_reader.pages])
                            except ImportError:
                                try:
                                    import PyPDF2
                                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(response.content))
                                    document_text = "\n".join([page.extract_text() or '' for page in pdf_reader.pages])
                                except ImportError:
                                    return "Error: PDF support requires 'pypdf' or 'PyPDF2' package. Install with: pip install pypdf"
                        except Exception as e:
                            return f"Error reading PDF: {str(e)}"
                    
                    elif 'word' in content_type or url.lower().endswith('.docx'):
                        try:
                            import io
                            from docx import Document
                            doc = Document(io.BytesIO(response.content))
                            document_text = "\n".join([para.text for para in doc.paragraphs])
                        except ImportError:
                            return "Error: Word document support requires 'python-docx' package. Install with: pip install python-docx"
                        except Exception as e:
                            return f"Error reading Word document: {str(e)}"
                    
                    elif 'presentation' in content_type or url.lower().endswith('.pptx'):
                        try:
                            import io
                            from pptx import Presentation
                            prs = Presentation(io.BytesIO(response.content))
                            texts = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        texts.append(shape.text)
                            document_text = "\n".join(texts)
                        except ImportError:
                            return "Error: PowerPoint support requires 'python-pptx' package. Install with: pip install python-pptx"
                        except Exception as e:
                            return f"Error reading PowerPoint: {str(e)}"
                    
                    else:
                        # Try as plain text
                        document_text = response.text
                        
                except Exception as e:
                    return f"Error fetching document from URL: {str(e)}"
            
            if not document_text:
                return "Error: No document content provided. Use 'url:URL' or 'content:TEXT'"
            
            # Truncate if too long (keep first 15000 chars for context)
            if len(document_text) > 15000:
                document_text = document_text[:15000] + "\n\n[Document truncated due to length...]"
            
            # Create prompt based on action
            if action == 'summarize':
                prompt = f"""Summarize the following document comprehensively. Include:
1. Main topic and purpose
2. Key sections/chapters overview
3. Main arguments or findings
4. Conclusions

Document:
{document_text}

Provide a well-structured summary:"""

            elif action == 'extract_key_points':
                prompt = f"""Extract the key points from this document. For each point:
- State the main idea clearly
- Note any supporting evidence
- Highlight important facts or figures

Document:
{document_text}

Key Points:"""

            elif action == 'answer_question':
                if not query:
                    return "Error: Please provide a question using 'query:YOUR_QUESTION'"
                prompt = f"""Based on the following document, answer this question: {query}

Document:
{document_text}

Answer the question thoroughly, citing specific parts of the document where relevant:"""

            elif action == 'find_quotes':
                if not query:
                    return "Error: Please provide a topic using 'query:TOPIC'"
                prompt = f"""Find relevant quotes and passages in this document about: {query}

Document:
{document_text}

List relevant quotes with their context:"""

            else:
                return f"Error: Unknown action '{action}'. Use: summarize, extract_key_points, answer_question, find_quotes"
            
            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"Document Analysis ({action}):\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in document reader: {e}")
            return f"Error analyzing document: {str(e)}"

    def _summarize_youtube(self, input_data: str) -> str:
        """YouTube Summarizer: Analyze videos for summaries, timestamps, insights"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            url = params.get('url', input_data.strip())
            action = params.get('action', 'full')
            
            if not url:
                return "Error: Please provide a YouTube URL. Format: 'url:https://www.youtube.com/watch?v=VIDEO_ID'"
            
            # Extract video ID
            video_id = None
            if 'youtube.com/watch?v=' in url:
                video_id = url.split('v=')[1].split('&')[0]
            elif 'youtu.be/' in url:
                video_id = url.split('youtu.be/')[1].split('?')[0]
            elif len(url) == 11:  # Just the ID
                video_id = url
            
            if not video_id:
                return "Error: Could not extract video ID from URL. Use format: https://www.youtube.com/watch?v=VIDEO_ID"
            
            # Get video transcript
            transcript_text = ""
            video_title = ""
            video_description = ""
            
            try:
                # Try to get transcript using youtube-transcript-api
                try:
                    from youtube_transcript_api import YouTubeTranscriptApi
                    transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                    
                    # Format transcript with timestamps
                    transcript_parts = []
                    for entry in transcript_list:
                        timestamp = int(entry['start'])
                        minutes = timestamp // 60
                        seconds = timestamp % 60
                        transcript_parts.append(f"[{minutes}:{seconds:02d}] {entry['text']}")
                    
                    transcript_text = "\n".join(transcript_parts)
                except ImportError:
                    self.logger.warning("youtube-transcript-api not installed, trying alternative method")
                except Exception as e:
                    self.logger.warning(f"Could not get transcript: {e}")
            except:
                pass
            
            # Try to get video metadata from oembed
            try:
                oembed_url = f"https://www.youtube.com/oembed?url=https://www.youtube.com/watch?v={video_id}&format=json"
                response = requests.get(oembed_url, timeout=10)
                if response.status_code == 200:
                    data = response.json()
                    video_title = data.get('title', '')
            except:
                pass
            
            if not transcript_text:
                # Fallback: provide instructions
                return f"""YouTube Video Analysis

Video ID: {video_id}
Title: {video_title or 'Unable to fetch'}
URL: https://www.youtube.com/watch?v={video_id}

Note: Could not retrieve video transcript automatically. 
To analyze this video, you can:
1. Install youtube-transcript-api: pip install youtube-transcript-api
2. Or provide the transcript manually using: 'content:TRANSCRIPT_TEXT,action:{action}'

The transcript is required to generate summaries, timestamps, insights, and action items."""
            
            # Truncate if too long
            if len(transcript_text) > 20000:
                transcript_text = transcript_text[:20000] + "\n\n[Transcript truncated...]"
            
            # Create prompt based on action
            if action == 'full':
                prompt = f"""Analyze this YouTube video transcript and provide:

1. SUMMARY (2-3 paragraphs)
2. KEY TIMESTAMPS (list the most important moments with their timestamps)
3. KEY INSIGHTS (main learnings and takeaways)
4. ACTION ITEMS (what viewers should do based on this video)

Video Title: {video_title}
Transcript:
{transcript_text}

Provide comprehensive analysis:"""

            elif action == 'summary':
                prompt = f"""Summarize this YouTube video:

Video Title: {video_title}
Transcript:
{transcript_text}

Provide a clear, comprehensive summary:"""

            elif action == 'timestamps':
                prompt = f"""Identify the key moments in this YouTube video with timestamps:

Video Title: {video_title}
Transcript:
{transcript_text}

List the most important moments with their timestamps and brief descriptions:"""

            elif action == 'insights':
                prompt = f"""Extract the key insights and learnings from this YouTube video:

Video Title: {video_title}
Transcript:
{transcript_text}

List the main insights, learnings, and important points:"""

            elif action == 'action_items':
                prompt = f"""What are the actionable takeaways from this YouTube video?

Video Title: {video_title}
Transcript:
{transcript_text}

List specific action items viewers can implement:"""

            else:
                return f"Error: Unknown action '{action}'. Use: full, summary, timestamps, insights, action_items"
            
            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"YouTube Video Analysis\nVideo: {video_title}\nURL: https://www.youtube.com/watch?v={video_id}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in YouTube summarizer: {e}")
            return f"Error analyzing YouTube video: {str(e)}"

    def _search_academic(self, input_data: str) -> str:
        """Academic Paper Search: Search and analyze research papers"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            query = params.get('query', input_data.strip())
            action = params.get('action', 'search')
            
            if not query:
                return "Error: Please provide a search query. Format: 'query:YOUR_SEARCH_QUERY'"
            
            papers = []
            
            # Search Semantic Scholar API (free, no auth required)
            try:
                self.logger.info(f"Searching Semantic Scholar for: {query}")
                ss_url = f"https://api.semanticscholar.org/graph/v1/paper/search"
                ss_params = {
                    'query': query,
                    'limit': 10,
                    'fields': 'title,authors,abstract,year,citationCount,url,openAccessPdf'
                }
                response = requests.get(ss_url, params=ss_params, timeout=15)
                
                if response.status_code == 200:
                    data = response.json()
                    for paper in data.get('data', []):
                        papers.append({
                            'title': paper.get('title', 'No title'),
                            'authors': ', '.join([a.get('name', '') for a in paper.get('authors', [])[:3]]),
                            'year': paper.get('year', 'N/A'),
                            'abstract': paper.get('abstract', 'No abstract available'),
                            'citations': paper.get('citationCount', 0),
                            'url': paper.get('url', ''),
                            'pdf': paper.get('openAccessPdf', {}).get('url', '') if paper.get('openAccessPdf') else '',
                            'source': 'Semantic Scholar'
                        })
            except Exception as e:
                self.logger.warning(f"Semantic Scholar search failed: {e}")
            
            # Also search arXiv (free, no auth required)
            if len(papers) < 5:
                try:
                    self.logger.info(f"Searching arXiv for: {query}")
                    arxiv_url = f"http://export.arxiv.org/api/query?search_query=all:{quote_plus(query)}&start=0&max_results=5"
                    response = requests.get(arxiv_url, timeout=15)
                    
                    if response.status_code == 200:
                        # Parse XML response
                        soup = BeautifulSoup(response.content, 'xml')
                        entries = soup.find_all('entry')
                        
                        for entry in entries:
                            title = entry.find('title')
                            summary = entry.find('summary')
                            authors = entry.find_all('author')
                            published = entry.find('published')
                            link = entry.find('id')
                            
                            papers.append({
                                'title': title.text.strip() if title else 'No title',
                                'authors': ', '.join([a.find('name').text for a in authors[:3]]) if authors else 'Unknown',
                                'year': published.text[:4] if published else 'N/A',
                                'abstract': summary.text.strip()[:500] if summary else 'No abstract',
                                'citations': 'N/A',
                                'url': link.text if link else '',
                                'pdf': link.text.replace('abs', 'pdf') + '.pdf' if link else '',
                                'source': 'arXiv'
                            })
                except Exception as e:
                    self.logger.warning(f"arXiv search failed: {e}")
            
            if not papers:
                return f"No academic papers found for: {query}. Try different keywords or check your internet connection."
            
            # Format results based on action
            if action == 'search':
                result = f"Academic Papers for: '{query}'\n\n"
                for i, paper in enumerate(papers[:10], 1):
                    result += f"{i}. {paper['title']}\n"
                    result += f"   Authors: {paper['authors']}\n"
                    result += f"   Year: {paper['year']} | Citations: {paper['citations']} | Source: {paper['source']}\n"
                    if paper['url']:
                        result += f"   URL: {paper['url']}\n"
                    if paper['pdf']:
                        result += f"   PDF: {paper['pdf']}\n"
                    result += f"   Abstract: {paper['abstract'][:200]}...\n\n"
                return result
            
            elif action in ['summarize', 'explain_simple', 'extract_methods']:
                # Use LLM to analyze papers
                papers_text = "\n\n".join([
                    f"Paper {i}: {p['title']}\nAuthors: {p['authors']}\nYear: {p['year']}\nAbstract: {p['abstract']}"
                    for i, p in enumerate(papers[:5], 1)
                ])
                
                if action == 'summarize':
                    prompt = f"""Summarize these academic papers on '{query}':

{papers_text}

Provide:
1. Overview of the research area
2. Key findings across papers
3. Common methodologies
4. Research gaps or future directions"""

                elif action == 'explain_simple':
                    prompt = f"""Explain these academic papers in simple terms (like explaining to a 5-year-old):

{papers_text}

Make it easy to understand:
1. What problem are they trying to solve?
2. How are they solving it?
3. What did they find out?
4. Why does it matter?"""

                elif action == 'extract_methods':
                    prompt = f"""Extract the methodologies, formulas, and techniques from these papers:

{papers_text}

List:
1. Research methods used
2. Key formulas or algorithms
3. Experimental setups
4. Evaluation metrics"""
                
                llm_caller = self._get_llm_caller()
                analysis = llm_caller.generate(prompt)
                
                result = f"Academic Paper Analysis: '{query}'\n\n"
                result += analysis + "\n\n"
                result += "--- Source Papers ---\n"
                for i, paper in enumerate(papers[:5], 1):
                    result += f"{i}. {paper['title']} ({paper['year']})\n   {paper['url']}\n"
                
                return result
            
            else:
                return f"Error: Unknown action '{action}'. Use: search, summarize, explain_simple, extract_methods"
            
        except Exception as e:
            self.logger.error(f"Error in academic search: {e}")
            return f"Error searching academic papers: {str(e)}"

    def _generate_mind_map(self, input_data: str) -> str:
        """Mind Map Generator: Create structured concept maps for any topic"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            topic = params.get('topic', input_data.strip())
            depth = int(params.get('depth', 2))
            depth = max(1, min(3, depth))  # Clamp between 1-3
            
            if not topic:
                return "Error: Please provide a topic. Format: 'topic:YOUR_TOPIC'"
            
            # Create prompt based on depth
            depth_instructions = {
                1: "Create a high-level overview with 4-6 main branches, each with 2-3 sub-concepts.",
                2: "Create a moderately detailed map with 5-7 main branches, each with 3-5 sub-concepts and their relationships.",
                3: "Create a deep, comprehensive map with 6-8 main branches, each with 4-6 sub-concepts, sub-sub-concepts, and detailed interconnections."
            }
            
            prompt = f"""Generate a comprehensive mind map for the topic: "{topic}"

{depth_instructions[depth]}

Structure your response as follows:

## CENTRAL CONCEPT
[Brief definition of the main topic]

## MAIN BRANCHES

### Branch 1: [Category Name]
- **Sub-concept 1.1**: [Brief explanation]
  - Detail: [Additional info]
- **Sub-concept 1.2**: [Brief explanation]
  - Detail: [Additional info]
- **Sub-concept 1.3**: [Brief explanation]

### Branch 2: [Category Name]
[Continue pattern...]

## RELATIONSHIPS & CONNECTIONS
- [Concept A] ↔ [Concept B]: [How they relate]
- [Concept C] → [Concept D]: [Causal relationship]
[List 5-10 key relationships between concepts]

## KEY TERMS GLOSSARY
- **Term 1**: Definition
- **Term 2**: Definition
[5-10 essential terms]

## LEARNING PATH
1. Start with: [concept]
2. Then explore: [concept]
3. Advanced: [concept]

Make it visual with clear hierarchies, use arrows (→, ↔, ←) for relationships, and ensure concepts flow logically."""

            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"🧠 MIND MAP: {topic}\n{'='*50}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in mind map generator: {e}")
            return f"Error generating mind map: {str(e)}"

    def _analyze_debate(self, input_data: str) -> str:
        """Debate Analyzer: Comprehensive pros/cons analysis for decisions"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            question = params.get('question', params.get('topic', input_data.strip()))
            
            if not question:
                return "Error: Please provide a question or topic. Format: 'question:Should I...?' or 'topic:Topic to analyze'"
            
            prompt = f"""Analyze the following decision/topic comprehensively: "{question}"

Provide a thorough analysis in the following structure:

## CONTEXT
[Brief overview of the decision and its significance]

## PROS (Advantages & Benefits)
1. **[Pro Title]**: [Detailed explanation]
   - Impact: [High/Medium/Low]
   - Timeframe: [Short-term/Long-term]
2. **[Pro Title]**: [Detailed explanation]
   - Impact: [High/Medium/Low]
   - Timeframe: [Short-term/Long-term]
[List 5-8 pros with details]

## CONS (Disadvantages & Drawbacks)
1. **[Con Title]**: [Detailed explanation]
   - Severity: [High/Medium/Low]
   - Mitigation: [How to reduce this risk]
2. **[Con Title]**: [Detailed explanation]
   - Severity: [High/Medium/Low]
   - Mitigation: [How to reduce this risk]
[List 5-8 cons with details]

## RISKS
1. **[Risk]**: [Description]
   - Probability: [High/Medium/Low]
   - Impact if occurs: [Description]
2. [Continue pattern...]
[List 4-6 key risks]

## OPPORTUNITY COSTS
What you give up by choosing this path:
1. [Opportunity cost 1]
2. [Opportunity cost 2]
[List 3-5 opportunity costs]

## LONG-TERM IMPACT ANALYSIS

### 1 Year
- Positive outcomes: [...]
- Potential challenges: [...]

### 5 Years
- Positive outcomes: [...]
- Potential challenges: [...]

### 10+ Years
- Positive outcomes: [...]
- Potential challenges: [...]

## KEY CONSIDERATIONS
- [Important factor to consider]
- [Another important factor]
[List 5-7 key things to think about]

## BALANCED RECOMMENDATION
[Provide a nuanced recommendation based on the analysis, acknowledging that the decision depends on individual circumstances. Include conditions under which each choice would be better.]

## DECISION FRAMEWORK
If [condition], then [recommendation].
If [different condition], then [different recommendation].
[Provide 3-4 conditional recommendations]"""

            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"⚖️ DECISION ANALYSIS: {question}\n{'='*50}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in debate analyzer: {e}")
            return f"Error analyzing decision: {str(e)}"

    def _first_principles_reasoning(self, input_data: str) -> str:
        """First Principles Reasoner: Break down problems to fundamental truths"""
        try:
            # Parse input
            params = {}
            for part in input_data.split(','):
                if ':' in part:
                    key, value = part.split(':', 1)
                    params[key.strip().lower()] = value.strip()
            
            problem = params.get('problem', params.get('idea', input_data.strip()))
            
            if not problem:
                return "Error: Please provide a problem or idea. Format: 'problem:YOUR_PROBLEM' or 'idea:YOUR_IDEA'"
            
            prompt = f"""Apply first-principles thinking to analyze: "{problem}"

First-principles thinking means breaking down complex problems into basic elements and then reassembling them from the ground up. This is how innovators like Elon Musk approach problems.

## STEP 1: DECOMPOSITION
Break down the problem into its fundamental components:

### Components Identified:
1. [Component]: [What it is and why it matters]
2. [Component]: [What it is and why it matters]
[List 5-8 fundamental components]

## STEP 2: QUESTION ASSUMPTIONS
What assumptions are commonly made about this problem?

### Assumptions to Challenge:
1. **Assumption**: [Common belief]
   - **Why it's assumed**: [Reason]
   - **Is it true?**: [Analysis]
   - **What if it's wrong?**: [Implication]
2. [Continue pattern...]
[List 5-7 assumptions]

## STEP 3: FUNDAMENTAL TRUTHS
What are the absolute, irreducible facts we know to be true?

### Core Truths:
1. **Truth**: [Fundamental fact]
   - **Why it's fundamental**: [Cannot be reduced further because...]
   - **Evidence**: [How we know this]
2. [Continue pattern...]
[List 5-8 fundamental truths - things that are objectively true and can't be broken down further]

## STEP 4: REBUILD FROM FIRST PRINCIPLES
Starting ONLY from the fundamental truths above, what solutions emerge?

### New Approach 1: [Name]
- **Built from truths**: [Which fundamental truths]
- **How it works**: [Description]
- **Why it's different**: [How this differs from conventional thinking]
- **Feasibility**: [High/Medium/Low]

### New Approach 2: [Name]
[Continue pattern...]

[Propose 2-4 novel approaches]

## STEP 5: NOVEL INSIGHTS
What new perspectives emerge from this analysis?

1. **Insight**: [Description]
   - **Implication**: [What this means]
2. [Continue pattern...]
[List 4-6 insights]

## STEP 6: PRACTICAL ACTION STEPS
Based on first-principles analysis, here's what to do:

### Immediate (This Week):
1. [Action]
2. [Action]

### Short-term (This Month):
1. [Action]
2. [Action]

### Long-term (This Quarter):
1. [Action]
2. [Action]

## KEY TAKEAWAY
[One paragraph summarizing the most important insight from this first-principles analysis]

## ELON MUSK WOULD SAY:
[A brief, provocative statement in the style of first-principles thinking that challenges conventional wisdom about this problem]"""

            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"🔬 FIRST PRINCIPLES ANALYSIS: {problem}\n{'='*50}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in first principles reasoner: {e}")
            return f"Error in first principles analysis: {str(e)}"

    def _generate_image(self, input_data: str) -> str:
        """Image Generator: Create images via Pollinations API - returns URL directly"""
        try:
            prompt = input_data.strip()
            
            if not prompt:
                return "Error: Please provide an image description."
            
            # URL encode the prompt
            encoded_prompt = quote_plus(prompt)
            
            # Build the API URL (Pollinations generates on-demand when URL is accessed)
            image_url = f"https://gen.pollinations.ai/image/{encoded_prompt}?model=flux"
            
            # Return just the URL - Pollinations generates when the link is clicked
            return image_url
            
        except Exception as e:
            self.logger.error(f"Error in image generator: {e}")
            return f"Error: {str(e)}"

    def _generate_story(self, input_data: str) -> str:
        """Story/Script Generator: Create stories and scripts for various media formats"""
        try:
            # Parse input
            params = {}
            
            if ':' in input_data and any(key in input_data.lower() for key in ['type:', 'premise:', 'genre:', 'tone:']):
                for part in input_data.split(','):
                    if ':' in part:
                        key, value = part.split(':', 1)
                        params[key.strip().lower()] = value.strip()
            else:
                # Treat as premise with default type
                params['premise'] = input_data.strip()
                params['type'] = 'short_story'
            
            story_type = params.get('type', 'short_story').lower()
            premise = params.get('premise', '')
            genre = params.get('genre', 'drama')
            tone = params.get('tone', 'engaging')
            length = params.get('length', 'medium')
            characters = params.get('characters', '')
            setting = params.get('setting', '')
            target = params.get('target', 'general audience')
            
            if not premise:
                return "Error: Please provide a premise. Format: 'type:comic,premise:Your story idea'"
            
            # Build format-specific prompt
            format_instructions = {
                'comic': """Create a comic book script with:
- Panel-by-panel breakdown (6-12 panels per page)
- Visual descriptions for each panel (what the artist should draw)
- Character dialogue in speech bubbles
- Sound effects (SFX) where appropriate
- Narrator captions if needed
- Page layout suggestions""",
                
                'game': """Create a game narrative with:
- Opening cinematic/hook
- Main quest structure with objectives
- Key dialogue trees for important conversations
- Character backstories and motivations
- World-building details
- Side quest ideas
- Key plot twists and reveals
- Multiple ending possibilities""",
                
                'movie': """Create a screenplay outline with:
- Logline (one-sentence summary)
- Three-act structure breakdown
- Scene-by-scene outline with sluglines
- Key dialogue moments
- Character introductions and arcs
- Visual motifs and themes
- Climax and resolution""",
                
                'ad': """Create an advertising script with:
- Hook (first 3 seconds to grab attention)
- Problem/pain point introduction
- Solution presentation
- Key benefits (3 max)
- Social proof element
- Call to action (CTA)
- Closing memorable line/tagline
- Visual/audio direction notes""",
                
                'short_story': """Create a short story with:
- Engaging opening hook
- Character introduction and development
- Rising action with tension
- Climax
- Resolution
- Thematic depth
- Vivid descriptions and dialogue""",
                
                'podcast': """Create a podcast script with:
- Cold open/teaser
- Intro music cue and host greeting
- Topic introduction
- Main content segments (2-3)
- Transition phrases
- Guest questions (if applicable)
- Audience engagement prompts
- Outro and call to action"""
            }
            
            format_guide = format_instructions.get(story_type, format_instructions['short_story'])
            
            prompt = f"""You are a professional storyteller and scriptwriter. Create a {story_type} based on the following:

**PREMISE:** {premise}
**GENRE:** {genre}
**TONE:** {tone}
**LENGTH:** {length}
**TARGET AUDIENCE:** {target}
{f'**KEY CHARACTERS:** {characters}' if characters else ''}
{f'**SETTING:** {setting}' if setting else ''}

**FORMAT REQUIREMENTS:**
{format_guide}

**DELIVERABLES:**

## 1. TITLE & LOGLINE
- Title: [Creative, memorable title]
- Logline: [One-sentence hook that captures the essence]

## 2. CHARACTER PROFILES
For each main character:
- Name
- Role in story
- Key traits
- Motivation
- Arc (how they change)

## 3. SETTING & WORLD
- Time period
- Location(s)
- Atmosphere/mood
- Key world rules (if fantasy/sci-fi)

## 4. STORY STRUCTURE
[Provide the full {story_type} content following the format requirements above]

## 5. KEY THEMES
- Primary theme
- Secondary themes
- How themes are explored

## 6. MEMORABLE MOMENTS
- 3-5 standout scenes/panels/beats that will resonate

Be creative, original, and true to the genre while making the content compelling and professional-quality."""

            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            type_emoji = {
                'comic': '📚',
                'game': '🎮',
                'movie': '🎬',
                'ad': '📺',
                'short_story': '📖',
                'podcast': '🎙️'
            }
            
            emoji = type_emoji.get(story_type, '✨')
            
            return f"{emoji} {story_type.upper()} SCRIPT: {premise[:50]}...\n{'='*50}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in story generator: {e}")
            return f"Error generating story: {str(e)}"

    def _plan_tasks(self, input_data: str) -> str:
        """Task Planner: Transform goals into actionable plans with steps, timelines, and checklists"""
        try:
            # Parse input
            params = {}
            
            if ':' in input_data and any(key in input_data.lower() for key in ['goal:', 'deadline:', 'budget:', 'hours']):
                for part in input_data.split(','):
                    if ':' in part:
                        key, value = part.split(':', 1)
                        params[key.strip().lower()] = value.strip()
            else:
                params['goal'] = input_data.strip()
            
            goal = params.get('goal', '')
            deadline = params.get('deadline', 'flexible')
            budget = params.get('budget', 'not specified')
            hours_per_week = params.get('hours_per_week', 'flexible')
            constraints = params.get('constraints', 'none specified')
            priority = params.get('priority', 'balanced')
            team_size = params.get('team_size', '1 (solo)')
            skill_level = params.get('skill_level', 'intermediate')
            
            if not goal:
                return "Error: Please provide a goal. Format: 'goal:Your goal here' or just describe your goal."
            
            prompt = f"""You are an expert project planner and productivity coach. Create a comprehensive, actionable plan for the following goal:

**GOAL:** {goal}

**CONTEXT:**
- Deadline: {deadline}
- Budget: {budget}
- Time available: {hours_per_week} hours per week
- Constraints: {constraints}
- Priority focus: {priority}
- Team size: {team_size}
- Skill level: {skill_level}

Create a detailed, practical plan with the following structure:

## 1. GOAL ANALYSIS

### Clarified Objective
[Restate the goal clearly and specifically]

### Success Criteria
How will you know when this is achieved? List 3-5 measurable outcomes:
- [ ] Criterion 1
- [ ] Criterion 2
- [ ] Criterion 3

### Scope Definition
- **In scope:** [What's included]
- **Out of scope:** [What's NOT included to keep focused]

## 2. PROJECT PHASES

Break the project into 3-5 major phases:

### Phase 1: [Name] (Week/Month X-Y)
**Milestone:** [What marks completion of this phase]
**Key deliverables:**
- Deliverable 1
- Deliverable 2

### Phase 2: [Name] (Week/Month X-Y)
[Continue pattern...]

## 3. DETAILED TASK BREAKDOWN

### Phase 1 Tasks:
| # | Task | Time Est. | Priority | Dependencies |
|---|------|-----------|----------|--------------|
| 1.1 | [Task] | [X hours] | High/Med/Low | None |
| 1.2 | [Task] | [X hours] | High/Med/Low | 1.1 |

### Phase 2 Tasks:
[Continue pattern for each phase]

## 4. TIMELINE & SCHEDULE

### Visual Timeline
```
Week 1-2:  [████████] Phase 1 - Foundation
Week 3-4:  [████████] Phase 2 - Development  
Week 5-6:  [████████] Phase 3 - Refinement
Week 7:    [████]     Phase 4 - Launch
```

### Key Milestones
| Date/Week | Milestone | Deliverable |
|-----------|-----------|-------------|
| Week 2 | [Milestone] | [Deliverable] |

## 5. MASTER CHECKLIST

Copy this checklist to track progress:

### Getting Started
- [ ] Review and understand full plan
- [ ] Set up necessary tools/accounts
- [ ] Block time in calendar

### Phase 1: [Name]
- [ ] Task 1.1
- [ ] Task 1.2
- [ ] Task 1.3
- [ ] ✓ Phase 1 Complete - Milestone achieved

### Phase 2: [Name]
[Continue for all phases...]

### Final Steps
- [ ] Review against success criteria
- [ ] Document lessons learned
- [ ] Celebrate completion! 🎉

## 6. RESOURCES NEEDED

### Tools & Software
- [Tool 1] - [Purpose]
- [Tool 2] - [Purpose]

### Skills Required
- [Skill 1] - [How to acquire if needed]
- [Skill 2] - [How to acquire if needed]

### People/Help
- [Role/Person] - [For what]

### Budget Allocation (if applicable)
| Category | Estimated Cost |
|----------|---------------|
| [Category] | $X |

## 7. RISK MANAGEMENT

| Risk | Likelihood | Impact | Mitigation Strategy |
|------|------------|--------|---------------------|
| [Risk 1] | High/Med/Low | High/Med/Low | [Strategy] |
| [Risk 2] | High/Med/Low | High/Med/Low | [Strategy] |

## 8. QUICK WINS 🚀

Start with these easy tasks to build momentum:
1. **[Quick Win 1]** - Can be done in <30 min
2. **[Quick Win 2]** - Can be done in <30 min
3. **[Quick Win 3]** - Can be done today

## 9. WEEKLY FOCUS GUIDE

**Week 1 Focus:** [Primary focus and 3 must-do tasks]
**Week 2 Focus:** [Primary focus and 3 must-do tasks]
[Continue as needed...]

## 10. MOTIVATION & ACCOUNTABILITY

### Why This Matters
[Brief reminder of the benefits of achieving this goal]

### Progress Tracking
- Review checklist: [Daily/Weekly]
- Milestone check-ins: [Schedule]
- Accountability: [How to stay on track]

Make this plan realistic, actionable, and motivating. Adjust complexity based on the goal size."""

            # Call LLM
            llm_caller = self._get_llm_caller()
            response = llm_caller.generate(prompt)
            
            return f"📋 TASK PLAN: {goal[:50]}...\n{'='*50}\n\n{response}"
            
        except Exception as e:
            self.logger.error(f"Error in task planner: {e}")
            return f"Error creating task plan: {str(e)}"

    def _multi_agent_process(self, input_data: str) -> str:
        """Multi-Agent Collaborator: Simulate multiple AI agents working together"""
        try:
            # Parse input
            params = {}
            
            if ':' in input_data and any(key in input_data.lower() for key in ['task:', 'depth:', 'focus:', 'style:']):
                for part in input_data.split(','):
                    if ':' in part:
                        key, value = part.split(':', 1)
                        params[key.strip().lower()] = value.strip()
            else:
                params['task'] = input_data.strip()
            
            task = params.get('task', '')
            depth = params.get('depth', 'standard')
            focus = params.get('focus', 'balanced')
            style = params.get('style', 'professional')
            iterations = params.get('iterations', '1')
            
            if not task:
                return "Error: Please provide a task. Format: 'task:Your task here'"
            
            # Get LLM caller
            llm_caller = self._get_llm_caller()
            
            results = []
            
            # === AGENT 1: RESEARCHER ===
            researcher_prompt = f"""You are the RESEARCHER agent in a multi-agent collaboration system.

**YOUR ROLE:** Gather information, explore the topic deeply, find relevant facts, data, and perspectives.

**TASK:** {task}
**DEPTH LEVEL:** {depth}
**STYLE:** {style}

**YOUR MISSION:**
1. Identify the key aspects of this task that need research
2. Explore multiple angles and perspectives
3. Gather relevant facts, statistics, examples, and expert opinions
4. Note any gaps in knowledge or areas needing clarification
5. Organize findings for the Writer agent

**OUTPUT FORMAT:**

## RESEARCH FINDINGS

### Key Topics Identified
1. [Topic 1]
2. [Topic 2]
3. [Topic 3]

### Core Facts & Information
- **[Aspect 1]:** [Findings]
- **[Aspect 2]:** [Findings]
- **[Aspect 3]:** [Findings]

### Multiple Perspectives
- **Perspective A:** [View and supporting points]
- **Perspective B:** [Contrasting view and supporting points]

### Relevant Examples/Case Studies
1. [Example 1]
2. [Example 2]

### Data & Statistics (if applicable)
- [Stat 1]
- [Stat 2]

### Key Questions to Address
1. [Question the content should answer]
2. [Question the content should answer]

### Notes for Writer
[Any important context or recommendations for the Writer agent]

Be thorough but focused. Quality over quantity."""

            researcher_response = llm_caller.generate(researcher_prompt)
            results.append(("🔍 RESEARCHER", researcher_response))
            
            # === AGENT 2: WRITER ===
            writer_prompt = f"""You are the WRITER agent in a multi-agent collaboration system.

**YOUR ROLE:** Create compelling content based on the Researcher's findings.

**ORIGINAL TASK:** {task}
**STYLE:** {style}

**RESEARCH PROVIDED BY RESEARCHER AGENT:**
{researcher_response}

**YOUR MISSION:**
1. Synthesize the research into well-structured content
2. Create engaging, clear, and purposeful writing
3. Address all key questions identified by the Researcher
4. Use examples and data effectively
5. Maintain the requested style throughout

**OUTPUT FORMAT:**

## INITIAL DRAFT

[Write the complete content here. This should be:
- Well-organized with clear sections
- Engaging and readable
- Comprehensive but not bloated
- True to the research findings
- Appropriate for the target audience]

### Structure Used:
- Introduction: [Brief note on approach]
- Body: [How organized]
- Conclusion: [What it achieves]

Make this draft as complete and polished as possible while staying open to critique."""

            writer_response = llm_caller.generate(writer_prompt)
            results.append(("✍️ WRITER", writer_response))
            
            # === AGENT 3: CRITIC ===
            critic_prompt = f"""You are the CRITIC agent in a multi-agent collaboration system.

**YOUR ROLE:** Evaluate the Writer's work, identify weaknesses, and suggest specific improvements.

**ORIGINAL TASK:** {task}

**RESEARCH (for reference):**
{researcher_response[:1500]}...

**WRITER'S DRAFT TO REVIEW:**
{writer_response}

**YOUR MISSION:**
1. Evaluate the content objectively and thoroughly
2. Identify strengths (to keep) and weaknesses (to fix)
3. Check for accuracy, clarity, and completeness
4. Suggest specific, actionable improvements
5. Be constructive but honest

**OUTPUT FORMAT:**

## CRITICAL REVIEW

### Overall Assessment
**Score:** [X/10]
**Summary:** [2-3 sentence overall evaluation]

### Strengths (Keep These)
1. ✅ [Strength 1 - what works well]
2. ✅ [Strength 2 - what works well]
3. ✅ [Strength 3 - what works well]

### Weaknesses (Need Improvement)
1. ⚠️ [Weakness 1]
   - **Problem:** [What's wrong]
   - **Impact:** [Why it matters]
   - **Fix:** [Specific suggestion]

2. ⚠️ [Weakness 2]
   - **Problem:** [What's wrong]
   - **Impact:** [Why it matters]
   - **Fix:** [Specific suggestion]

3. ⚠️ [Weakness 3]
   - **Problem:** [What's wrong]
   - **Impact:** [Why it matters]
   - **Fix:** [Specific suggestion]

### Missing Elements
- [What should be added]
- [What should be added]

### Accuracy Check
- [Any factual issues or unsupported claims]

### Clarity Issues
- [Any confusing passages and how to clarify]

### Priority Improvements (Top 3)
1. 🔴 [Most important fix]
2. 🟡 [Second priority]
3. 🟢 [Third priority]

Be tough but fair. The goal is to make this excellent."""

            critic_response = llm_caller.generate(critic_prompt)
            results.append(("🔎 CRITIC", critic_response))
            
            # === AGENT 4: WRITER REVISION ===
            revision_prompt = f"""You are the WRITER agent again, now revising based on the Critic's feedback.

**ORIGINAL TASK:** {task}

**YOUR ORIGINAL DRAFT:**
{writer_response}

**CRITIC'S FEEDBACK:**
{critic_response}

**YOUR MISSION:**
1. Address ALL the Critic's priority improvements
2. Incorporate the suggested fixes
3. Keep the identified strengths
4. Add any missing elements
5. Improve clarity where noted

**OUTPUT FORMAT:**

## REVISED DRAFT

[Write the complete revised content here, incorporating all feedback]

### Changes Made:
1. [Change 1 - addressing which critique]
2. [Change 2 - addressing which critique]
3. [Change 3 - addressing which critique]

Make this the best possible version."""

            revision_response = llm_caller.generate(revision_prompt)
            results.append(("✍️ WRITER (Revised)", revision_response))
            
            # === AGENT 5: SUMMARIZER ===
            summarizer_prompt = f"""You are the SUMMARIZER agent in a multi-agent collaboration system.

**YOUR ROLE:** Create the final polished output with key takeaways.

**ORIGINAL TASK:** {task}

**FINAL REVISED CONTENT:**
{revision_response}

**YOUR MISSION:**
1. Polish the final content
2. Create an executive summary
3. Extract key takeaways
4. Provide the final deliverable

**OUTPUT FORMAT:**

## FINAL OUTPUT

### Executive Summary
[3-5 sentence summary of the main content and conclusions]

### Key Takeaways
1. 💡 [Takeaway 1]
2. 💡 [Takeaway 2]
3. 💡 [Takeaway 3]
4. 💡 [Takeaway 4]
5. 💡 [Takeaway 5]

### Final Polished Content
[The complete, polished final version - ready for use]

### Action Items (if applicable)
- [ ] [Action 1]
- [ ] [Action 2]
- [ ] [Action 3]

### Quality Score
**Final Rating:** [X/10]
**Confidence Level:** [High/Medium/Low]

This represents the collaborative output of the multi-agent team."""

            summarizer_response = llm_caller.generate(summarizer_prompt)
            results.append(("📝 SUMMARIZER", summarizer_response))
            
            # Compile final output
            output = f"🤖 MULTI-AGENT COLLABORATION: {task[:50]}...\n{'='*60}\n\n"
            output += "This task was processed by 4 specialized AI agents working together:\n"
            output += "Researcher → Writer → Critic → Writer (Revision) → Summarizer\n\n"
            output += "="*60 + "\n\n"
            
            for agent_name, agent_response in results:
                output += f"## {agent_name}\n"
                output += "-"*40 + "\n\n"
                output += agent_response + "\n\n"
                output += "="*60 + "\n\n"
            
            return output
            
        except Exception as e:
            self.logger.error(f"Error in multi-agent process: {e}")
            return f"Error in multi-agent collaboration: {str(e)}"